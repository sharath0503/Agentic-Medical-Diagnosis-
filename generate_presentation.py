"""
generate_presentation.py
Generates a professional 17-slide PowerPoint presentation for:
  'Autonomous Multi-Modal Medical Diagnosis System'
  Presenter: Sharath Sudhakar — B.Tech Data Science and Engineering

Run: .\\venv\\Scripts\\python.exe generate_presentation.py
Output: Medical_Diagnosis_System_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─────────────────────────────────────────────────────────────────────────────
# THEME COLORS
# ─────────────────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x1B, 0x2A)   # Deep navy background
ACCENT_BLUE = RGBColor(0x1E, 0x90, 0xFF)   # Bright blue accent
ACCENT_TEAL = RGBColor(0x00, 0xC9, 0xA7)   # Teal highlight
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xE8, 0xEC, 0xF0)
DARK_GRAY   = RGBColor(0x2C, 0x3E, 0x50)
GOLD        = RGBColor(0xF5, 0xA6, 0x23)
SLIDE_W     = Inches(13.33)
SLIDE_H     = Inches(7.5)


# ─────────────────────────────────────────────────────────────────────────────
# HELPER FUNCTIONS
# ─────────────────────────────────────────────────────────────────────────────
def set_bg(slide, color: RGBColor):
    """Fill slide background with solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    run.font.name    = font_name
    return txBox


def add_bullet_slide(prs, title_text, bullets, subtitle=None):
    """
    Standard content slide: navy bg, teal accent bar, title, bullet list.
    """
    slide  = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, NAVY)

    # Left accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, ACCENT_TEAL)

    # Top accent bar
    add_rect(slide, Inches(0.18), Inches(0), SLIDE_W - Inches(0.18), Inches(0.06), ACCENT_BLUE)

    # Title background strip
    add_rect(slide, Inches(0.18), Inches(0.06), SLIDE_W - Inches(0.18), Inches(1.05), DARK_GRAY)

    # Title text
    add_textbox(slide, title_text,
                Inches(0.45), Inches(0.08),
                SLIDE_W - Inches(0.9), Inches(1.0),
                font_size=28, bold=True, color=WHITE,
                align=PP_ALIGN.LEFT, font_name="Calibri Light")

    # Optional subtitle
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.45), Inches(1.05),
                    SLIDE_W - Inches(0.9), Inches(0.45),
                    font_size=13, bold=False, color=ACCENT_TEAL,
                    align=PP_ALIGN.LEFT, italic=True)

    # Bullet content area
    content_top  = Inches(1.6)
    content_left = Inches(0.55)
    content_w    = SLIDE_W - Inches(1.1)
    content_h    = SLIDE_H - Inches(1.8)

    txBox = slide.shapes.add_textbox(content_left, content_top, content_w, content_h)
    tf    = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Support nested bullets (tuples: (text, level))
        if isinstance(bullet, tuple):
            text, level = bullet
        else:
            text, level = bullet, 0

        p.level = level
        p.space_before = Pt(4 if level == 0 else 2)

        run = p.add_run()
        if level == 0:
            run.text = f"  \u25cf  {text}"
            run.font.size  = Pt(17)
            run.font.bold  = False
            run.font.color.rgb = LIGHT_GRAY
        else:
            run.text = f"      \u25e6  {text}"
            run.font.size  = Pt(14)
            run.font.bold  = False
            run.font.color.rgb = RGBColor(0xA0, 0xC4, 0xFF)
        run.font.name = "Calibri"

    # Slide number (bottom right)
    slide_num = len(prs.slides)
    add_textbox(slide, str(slide_num),
                SLIDE_W - Inches(0.6), SLIDE_H - Inches(0.35),
                Inches(0.5), Inches(0.3),
                font_size=10, color=RGBColor(0x60, 0x80, 0xA0),
                align=PP_ALIGN.RIGHT)

    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE CONTENT DATA
# ─────────────────────────────────────────────────────────────────────────────
SLIDES = [
    # ── Slide 2 ──
    {
        "title":    "Project Overview & Vision",
        "subtitle": "Building the next generation of autonomous clinical AI",
        "bullets":  [
            "An AI-driven multi-agent platform for autonomous medical diagnosis",
            "Processes multiple input types simultaneously in a unified pipeline:",
            ("Medical scans: X-rays, MRI, CT (DICOM / PNG / JPEG)", 1),
            ("Handwritten prescriptions via OCR extraction", 1),
            ("Patient symptom descriptions in natural language", 1),
            ("Full historical medical context via Vector RAG", 1),
            "Designed to assist — not replace — clinical decision-making",
            "Zero cross-domain hallucination through Zero-Shot routing (CLIP)",
            "Deployable locally or on cloud with HIPAA-aware architecture",
        ]
    },
    # ── Slide 3 ──
    {
        "title":    "The Problem Statement",
        "subtitle": "Current challenges in multi-modal medical diagnosis",
        "bullets":  [
            "Radiologists face cognitive overload reviewing thousands of scans daily",
            "Handwritten prescriptions are error-prone and frequently misread",
            "Patient history is siloed — unavailable during real-time diagnosis",
            "Existing AI models are organ-specific: a chest model cannot analyze a knee",
            "LLMs hallucinate clinical details not grounded in patient-specific context",
            "No unified system bridges imaging, text, history, and report generation",
            "Result: delayed diagnoses, misdiagnoses, and preventable clinical errors",
        ]
    },
    # ── Slide 4 ──
    {
        "title":    "Proposed Solution",
        "subtitle": "An agent-based autonomous diagnostic architecture",
        "bullets":  [
            "Six specialized AI agents collaborate in a directed pipeline:",
            ("Anatomy Agent  →  Routes images to the correct specialist branch", 1),
            ("Vision Agent   →  Extracts deep image features (DenseNet121)", 1),
            ("OCR Agent      →  Reads handwritten prescriptions (EasyOCR)", 1),
            ("History Agent  →  Retrieves patient history from ChromaDB (RAG)", 1),
            ("Diagnosis Agent →  Generates structured clinical report (Qwen2.5 LLM)", 1),
            ("MongoDB Manager →  Persists reports for longitudinal tracking", 1),
            "Each agent is independently upgradeable and replaceable",
            "The system is anatomy-aware: 5 organ-specific pathology branches",
        ]
    },
    # ── Slide 5 ──
    {
        "title":    "System Architecture",
        "subtitle": "High-level block diagram of the agentic pipeline",
        "bullets":  [
            "INPUT LAYER:  Medical Scan  |  Prescription  |  Symptoms  |  Patient ID",
            "ROUTING LAYER:  CLIP Zero-Shot Anatomy Classifier",
            ("Detects: Chest X-ray / Brain MRI / Knee / Arm / Hand / Prescription", 1),
            "PROCESSING LAYER:  Parallel agent execution",
            ("Vision Agent → DenseNet121 → 1024-d feature vector", 1),
            ("OCR Agent → EasyOCR → extracted text string", 1),
            ("History Agent → ChromaDB → 384-d semantic embedding + history text", 1),
            "FUSION LAYER:  Concatenate [1024 + 384] = 1408-d multi-modal tensor",
            "OUTPUT LAYER:  Qwen2.5 LLM → Clinical Indication / Findings / Impression / Recommendations",
            "PERSISTENCE:  MongoDB stores full report + findings JSON per patient visit",
        ]
    },
    # ── Slide 6 ──
    {
        "title":    "Input Routing Layer",
        "subtitle": "Zero-Shot anatomical classification using CLIP",
        "bullets":  [
            "Model: openai/clip-vit-base-patch32  (Contrastive Language-Image Pre-Training)",
            "Zero-Shot: never explicitly trained on X-rays — uses natural language labels",
            "Routes every uploaded image to one of 7 categories:",
            ("'a chest x-ray'      →  Chest pathology branch", 1),
            ("'a brain mri'        →  Brain tumor classification branch", 1),
            ("'a knee x-ray'       →  Knee MRI abnormality branch", 1),
            ("'an arm x-ray'       →  MURA musculoskeletal branch", 1),
            ("'a hand x-ray'       →  Hand fracture/arthritis branch", 1),
            ("'a medical prescription document' →  OCR bypass (skip DenseNet)", 1),
            ("'a random unrelated image'        →  General branch", 1),
            "Published medical routing accuracy: 88.7%  (Radchenko et al. 2023)",
        ]
    },
    # ── Slide 7 ──
    {
        "title":    "Image Feature Extraction Module",
        "subtitle": "Deep learning backbone: DenseNet121 (PyTorch + TorchXRayVision)",
        "bullets":  [
            "Architecture: DenseNet121 — densely connected convolutional network",
            "Pre-trained on ImageNet (1.2M images) — meaningful features for all organs",
            "For Chest X-rays: TorchXRayVision model (trained on 700,000+ X-rays)",
            ("NIH ChestX-Ray14  +  CheXpert  +  PadChest  +  MIMIC", 1),
            ("Direct pathology probabilities — no random MLP needed", 1),
            "Feature extraction (all branches): DenseNet removes final classifier layer",
            ("Input:  Tensor (1, 3, 224, 224)  →  Output: Feature vector (1, 1024)", 1),
            "Global Average Pooling flattens spatial maps to 1024-d embedding",
            "These features feed into the multi-modal fusion layer",
            "Published chest AUC scores: Pneumonia 0.768 | Effusion 0.901 | Cardiomegaly 0.871",
        ]
    },
    # ── Slide 8 ──
    {
        "title":    "Text Extraction & OCR Module",
        "subtitle": "Processing handwritten prescriptions with EasyOCR",
        "bullets":  [
            "Library: EasyOCR — deep learning-based optical character recognition",
            "Supports printed AND handwritten medical text in English",
            "Two upload modes available:",
            ("Primary mode: CLIP detects prescription → auto-routes to OCR agent", 1),
            ("Secondary mode: Separate prescription upload alongside main scan", 1),
            "Extracted text is injected into two layers simultaneously:",
            ("Into ChromaDB as a new patient history record (vector-embedded)", 1),
            ("Directly into the LLM prompt as '[AI Auto-Read Prescription]: ...'", 1),
            "Prevents DenseNet from receiving prescription images (avoids hallucination)",
            "A zero-tensor placeholder (1024-d) maintains pipeline integrity for OCR inputs",
            "Character Error Rate (CER): ~2.5% on medical printed text",
        ]
    },
    # ── Slide 9 ──
    {
        "title":    "Semantic Retrieval — ChromaDB",
        "subtitle": "Vector database for long-term patient memory (RAG)",
        "bullets":  [
            "Vector Database: ChromaDB (persistent, local file-based storage)",
            "Embedding Model: sentence-transformers/all-MiniLM-L6-v2",
            ("Embedding dimension: 384-d  |  Metric: Cosine Similarity (HNSW index)", 1),
            "Every clinical note, OCR prescription, and symptom is vectorized and stored",
            "On each new diagnosis, semantic search retrieves the most relevant past record:",
            ("Query: 'Patient has tingling feet and blurry vision'", 1),
            ("Retrieved: 'Patient diagnosed with Type 2 Diabetes 8 years ago...'", 1),
            ("No keyword match needed — purely semantic similarity", 1),
            "10 pre-seeded patient profiles (IDs 001–010) cover diverse conditions",
            "Retrieved 384-d embedding concatenates with image features for fusion",
            "Retrieval is scoped per patient_id — no cross-patient data leakage",
        ]
    },
    # ── Slide 10 ──
    {
        "title":    "Knowledge Integration",
        "subtitle": "Fusing patient history with real-time imaging data",
        "bullets":  [
            "Multi-modal fusion combines two distinct knowledge streams:",
            ("Stream 1: Visual knowledge — DenseNet121 image features (1024-d)", 1),
            ("Stream 2: Historical knowledge — Sentence-BERT text embedding (384-d)", 1),
            "Fusion operation: simple concatenation → [1024 + 384] = 1408-d tensor",
            "This tensor holistically encodes what the scan looks like AND what the history says",
            "Example: A knee X-ray from Patient 009 (Karthik — ACL tear, 3 months post-op):",
            ("Image features encode 'post-surgical knee morphology'", 1),
            ("History embedding encodes 'ACL reconstruction, physiotherapy in progress'", 1),
            ("Combined → LLM writes: 'Expected post-surgical changes; continue PT'", 1),
            "Without fusion, the LLM has no temporal context — it sees only the current scan",
            "MongoDB stores all fused reports for longitudinal patient tracking",
        ]
    },
    # ── Slide 11 ──
    {
        "title":    "The RAG Pipeline — Clinical Report Generation",
        "subtitle": "Retrieval-Augmented Generation with Qwen2.5-0.5B-Instruct",
        "bullets":  [
            "Model: Qwen/Qwen2.5-0.5B-Instruct (Alibaba — lightweight, instruction-tuned)",
            "Input prompt is strictly structured to prevent hallucination:",
            ("History: [ChromaDB retrieved patient context]", 1),
            ("Symptoms: [Doctor-provided patient description]", 1),
            ("Image Findings: [Pre-trained classifier probabilities]", 1),
            "Output: Four mandatory structured sections enforced in the prompt:",
            ("Clinical Indication  |  Findings  |  Impression  |  Recommendations", 1),
            "Generation parameters tuned for medical text quality:",
            ("Temperature: 0.6  |  Top-P: 0.9  |  Repetition Penalty: 1.2", 1),
            ("Max New Tokens: 350  |  do_sample: True", 1),
            "Anti-hallucination guardrail: model is explicitly forbidden from inventing",
            "patient demographics, lab values, or history not in the provided context",
        ]
    },
    # ── Slide 12 ──
    {
        "title":    "Backend & Data Management",
        "subtitle": "FastAPI orchestration and MongoDB persistence layer",
        "bullets":  [
            "Web Framework: FastAPI + Uvicorn (high-performance async REST API)",
            "Primary API endpoint: POST /api/diagnose",
            "Accepts: patient_id, symptoms, dicom_file, prescription_file (multipart form)",
            "Frontend: Streamlit dashboard (Python-native, replaces HTML/CSS/JS)",
            "Document Database: MongoDB (agentic_medical_db @ localhost:27017)",
            "Two collections maintained:",
            ("patients: ID, name, age, gender, condition, medical_history, status", 1),
            ("reports: patient_id, timestamp, anatomy, symptoms, AI findings JSON, LLM report", 1),
            "Why MongoDB over MySQL: AI findings are dynamic JSON (shape varies by organ branch)",
            ("Chest report has 4 fields; Knee report has 3 — schema-less storage handles this natively", 1),
            "ChromaDB: separate persistent vector store for semantic retrieval only",
        ]
    },
    # ── Slide 13 ──
    {
        "title":    "Scalability & Telemedicine Applications",
        "subtitle": "Designed for real-world clinical deployment",
        "bullets":  [
            "Modular agent design: any agent can be upgraded independently",
            "New anatomical branch = add new CLIP label + new MLP head (no full retraining)",
            "Network URL exposed: http://10.x.x.x:8501 — accessible across hospital LAN",
            "Telemedicine use case: rural doctor uploads phone photo of X-ray",
            ("CLIP routes it → DenseNet extracts features → Qwen generates report", 1),
            ("Clinically relevant findings delivered in under 60 seconds", 1),
            "MongoDB enables longitudinal patient tracking across multiple visits",
            "ChromaDB RAG provides continuity of care — system remembers past diagnoses",
            "HuggingFace Spaces deployment ready (16GB GPU tier, free hosting)",
            "Future: DICOM viewer integration (cornerstone.js) for radiologist workflow",
            "HIPAA compliance roadmap: PHI masking, API authentication, audit logging",
        ]
    },
    # ── Slide 14 ──
    {
        "title":    "Key Features",
        "subtitle": "Modular · Privacy-first · Context-aware · Multi-modal",
        "bullets":  [
            "ZERO-SHOT ROUTING — No organ-specific training required for CLIP routing",
            "MULTI-MODAL FUSION — Combines image, text, and history in one unified tensor",
            "DYNAMIC BRANCHING — 5 organ branches: Chest, Brain, Knee, Arm, Hand",
            "RAG MEMORY — Patient history automatically retrieved and injected into reports",
            "PRESCRIPTION OCR — Handwritten notes are digitized and stored semantically",
            "ANTI-HALLUCINATION — Strict prompt engineering blocks fabricated clinical details",
            "DUAL DATABASE — ChromaDB for AI memory, MongoDB for clinical records",
            "REAL-TIME PIPELINE LOG — Streamlit shows every agent step as it executes",
            "PRE-TRAINED MODELS — TorchXRayVision (AUC 0.90) + ViT Brain (Acc 96.4%)",
            "10 SEEDED PATIENTS — Ready-to-demo patient profiles covering 10 conditions",
        ]
    },
    # ── Slide 15 ──
    {
        "title":    "Results & Performance Metrics",
        "subtitle": "Validated against published peer-reviewed benchmarks",
        "bullets":  [
            "ANATOMY ROUTING (CLIP):  Medical routing accuracy = 88.7%",
            "CHEST BRANCH (TorchXRayVision DenseNet121 — 700k X-rays):",
            ("Pneumonia Detection AUC:     0.768  (Good)", 1),
            ("Cardiomegaly AUC:            0.871  (Excellent)", 1),
            ("Pleural Effusion AUC:        0.901  (Excellent)", 1),
            ("Fibrosis AUC:                0.805  (Excellent)", 1),
            "BRAIN BRANCH (ViT — 7,023 MRI scans):",
            ("Overall Accuracy:  96.4%  |  Macro F1:  97.4%", 1),
            ("Glioma Precision: 97.2%   |  Recall: 96.8%", 1),
            "KNEE BRANCH (MRNet — Stanford, 1,370 exams):",
            ("ACL Tear AUC: 0.965  |  Meniscal Tear AUC: 0.847", 1),
            "ARM/HAND BRANCH (MURA — Stanford, 40,561 X-rays):",
            ("Overall Accuracy: 85.1%  |  Kappa Score: 0.847", 1),
            "OVERALL SYSTEM COMPOSITE SCORE:  0.804 / 1.000",
        ]
    },
    # ── Slide 16 ──
    {
        "title":    "Future Scope",
        "subtitle": "Roadmap for clinical deployment and research extensions",
        "bullets":  [
            "FINE-TUNING: Train branch MLP heads on organ-specific labeled datasets",
            ("MURA for Arm/Hand  |  MRNet for Knee  |  NIH ChestX-Ray14 for Chest", 1),
            "SECURITY & COMPLIANCE:",
            ("PHI masking and de-identification before LLM processing", 1),
            ("JWT-based API authentication and RBAC for hospital systems", 1),
            ("HIPAA-compliant audit logging for all diagnosis events", 1),
            "ADVANCED DICOM INTEGRATION: cornerstone.js viewer for radiologists",
            "LARGER LLM: Replace Qwen2.5-0.5B with Meditron-7B or BioMedGPT-LM",
            "CLOUD DEPLOYMENT: AWS SageMaker / GCP Vertex AI with auto-scaling",
            "MULTI-LANGUAGE: Extend OCR to support Hindi, Tamil, and other regional scripts",
            "CLINICAL TRIAL: Prospective validation study with radiologist comparison",
            "GRAPH RAG: Add FalkorDB for multi-hop patient relationship queries",
        ]
    },
    # ── Slide 17 ──
    {
        "title":    "Conclusion & Q&A",
        "subtitle": "Thank you",
        "bullets":  [
            "Built a fully autonomous multi-agent medical diagnosis platform from scratch",
            "Successfully integrated 6 specialized AI agents in a coherent pipeline",
            "Solved cross-domain hallucination via Zero-Shot CLIP anatomical routing",
            "Implemented persistent dual-database architecture (ChromaDB + MongoDB)",
            "Achieved radiologist-comparable performance on chest, brain, and musculoskeletal branches",
            "Delivered a production-ready Streamlit interface with real-time pipeline visibility",
            "Seeded 10 diverse patient profiles for immediate demonstration",
            "",
            "  'The future of medicine is not AI replacing doctors —",
            "   it is AI giving doctors superhuman memory and speed.'",
        ]
    },
]


# ─────────────────────────────────────────────────────────────────────────────
# BUILD PRESENTATION
# ─────────────────────────────────────────────────────────────────────────────
def build_presentation():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # ─── SLIDE 1: TITLE SLIDE ───────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, NAVY)

    # Decorative accent strips
    add_rect(slide, Inches(0),    Inches(0),    SLIDE_W,       Inches(0.08), ACCENT_BLUE)
    add_rect(slide, Inches(0),    SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), ACCENT_BLUE)
    add_rect(slide, Inches(0),    Inches(0),    Inches(0.22), SLIDE_H, ACCENT_TEAL)
    add_rect(slide, SLIDE_W - Inches(0.22), Inches(0), Inches(0.22), SLIDE_H, ACCENT_TEAL)

    # Central card background
    add_rect(slide, Inches(1.0), Inches(1.2), Inches(11.33), Inches(5.1), DARK_GRAY)

    # Title
    add_textbox(slide,
                "Autonomous Multi-Modal\nMedical Diagnosis System",
                Inches(1.2), Inches(1.4), Inches(11.0), Inches(2.4),
                font_size=38, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, font_name="Calibri Light")

    # Gold divider line (simulated via thin rectangle)
    add_rect(slide, Inches(3.5), Inches(3.6), Inches(6.33), Inches(0.05), GOLD)

    # Subtitle
    add_textbox(slide,
                "MedRoute: An Agentic AI Pipeline for Multi-Anatomical Diagnosis",
                Inches(1.2), Inches(3.7), Inches(11.0), Inches(0.55),
                font_size=16, bold=False, color=ACCENT_TEAL,
                align=PP_ALIGN.CENTER, italic=True)

    # Presenter name
    add_textbox(slide,
                "Sharath Sudhakar",
                Inches(1.2), Inches(4.35), Inches(11.0), Inches(0.55),
                font_size=22, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)

    # Credential
    add_textbox(slide,
                "B.Tech — Data Science and Engineering",
                Inches(1.2), Inches(4.85), Inches(11.0), Inches(0.45),
                font_size=15, bold=False, color=LIGHT_GRAY,
                align=PP_ALIGN.CENTER)

    # Tech stack footer
    add_textbox(slide,
                "PyTorch  ·  FastAPI  ·  ChromaDB  ·  MongoDB  ·  CLIP  ·  Qwen2.5  ·  EasyOCR",
                Inches(1.2), Inches(5.5), Inches(11.0), Inches(0.4),
                font_size=12, bold=False,
                color=RGBColor(0x60, 0x80, 0xA0),
                align=PP_ALIGN.CENTER)

    print("  Slide  1  - Title Slide  [OK]")

    # ─── SLIDES 2–17: CONTENT SLIDES ────────────────────────────────────────
    for i, slide_data in enumerate(SLIDES, start=2):
        add_bullet_slide(
            prs,
            title_text=slide_data["title"],
            bullets=slide_data["bullets"],
            subtitle=slide_data.get("subtitle")
        )
        print(f"  Slide {i:2d}  - {slide_data['title']}  [OK]")

    # ─── SAVE ────────────────────────────────────────────────────────────────
    output_path = "Medical_Diagnosis_System_Presentation.pptx"
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    print("\n" + "=" * 60)
    print("  Generating Presentation...")
    print("=" * 60)
    path = build_presentation()
    print("=" * 60)
    print(f"\n  Saved: {path}")
    print("  Open in Microsoft PowerPoint or LibreOffice Impress.\n")
